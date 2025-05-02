from .BaseDocumentParser import BaseDocumentParser
from pptx import Presentation


class PPTXParser(BaseDocumentParser):
    def extract_text(self, file_path):
        prs = Presentation(file_path)
        text = []

        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text.append(shape.text)
        
        return "\n".join(text)